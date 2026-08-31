from pptx import Presentation
from src.document_parser.base_parser import BaseDocumentParser

class PptParser(BaseDocumentParser):
    def parse(self, file_path: str) -> str:
        prs = Presentation(file_path)
        all_text = []
        # 遍历每一页幻灯片
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                # 提取文本框文字
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        slide_text.append(txt)
            if slide_text:
                all_text.append(f"【第{slide_idx}页】\n" + "\n".join(slide_text))
        return "\n\n".join(all_text)